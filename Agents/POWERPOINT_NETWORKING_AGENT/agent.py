"""
POWERPOINT_NETWORKING_AGENT: Networking Strategy PowerPoint Generator

Creates professional networking strategy presentations with:
- Proper slide titles
- Exactly 3 bullets per content slide
- Clean spacing and typography
"""

import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any
import re
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

env_path = Path(__file__).parent.parent.parent / "Secrets" / ".env"
load_dotenv(env_path)


class SlideStyle:
    """Professional styling constants."""
    PRIMARY = RGBColor(0, 51, 102)
    SECONDARY = RGBColor(0, 102, 153)
    ACCENT = RGBColor(0, 153, 204)
    TEXT_DARK = RGBColor(51, 51, 51)
    TEXT_LIGHT = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)

    MARGIN_LEFT = Inches(0.75)
    MAX_BULLET_CHARS = 120


class PowerPointNetworkingAgent:
    """Agent for converting networking strategy to PowerPoint."""

    def __init__(self, api_key: str = None):
        self.openai_api_key = api_key or os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError("OpenAI API key not provided")

        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.3,
            openai_api_key=self.openai_api_key
        )

        self.base_path = Path(__file__).parent.parent.parent
        self.outputs_path = self.base_path / "Outputs"
        self.outputs_path.mkdir(exist_ok=True)
        self.style = SlideStyle()

    def load_strategy(self, strategy_file: Path = None, sector: str = None) -> str:
        """Load the networking strategy markdown file."""
        if strategy_file and strategy_file.exists():
            with open(strategy_file, 'r', encoding='utf-8') as f:
                return f.read()

        if sector:
            sector_slug = sector.lower().replace(" ", "_")
            pattern = f"networking_strategy_{sector_slug}_*.md"
            matching_files = list(self.outputs_path.glob(pattern))

            if matching_files:
                latest_file = max(matching_files, key=lambda p: p.stat().st_mtime)
                print(f"  Loading strategy: {latest_file.name}")
                with open(latest_file, 'r', encoding='utf-8') as f:
                    return f.read()

        raise ValueError("No networking strategy file found.")

    def parse_strategy(self, markdown: str) -> Dict[str, Any]:
        """Parse the strategy markdown into structured sections."""
        data = {
            'sector': '',
            'region': 'US',
            'sections': []
        }

        region_match = re.search(r'\*\*Region\*\*:\s*(.+?)$', markdown, re.MULTILINE)
        if region_match:
            data['region'] = region_match.group(1).strip()

        sector_match = re.search(r'^## (.+?)$', markdown, re.MULTILINE)
        if sector_match:
            data['sector'] = sector_match.group(1).strip()

        section_pattern = r'## (\d+)\. (.+?)\n(.*?)(?=\n## \d|\n---\s*\n## Methodology|$)'
        for match in re.finditer(section_pattern, markdown, re.DOTALL):
            section_num = match.group(1)
            section_title = match.group(2).strip()
            section_content = match.group(3).strip()

            data['sections'].append({
                'number': section_num,
                'title': section_title,
                'content': section_content
            })

        return data

    def condense_to_three_bullets(self, content: str, section_title: str) -> List[str]:
        """Ensure exactly 3 concise bullets."""
        system_prompt = """Convert the content into EXACTLY 3 bullet points.

Rules:
- Each bullet must be under 100 characters
- Keep key data points
- No sub-bullets
- Return ONLY 3 lines starting with '- '"""

        user_prompt = f"""Section: {section_title}

Content:
{content}

Convert to exactly 3 concise bullet points:"""

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]

        response = self.llm.invoke(messages)

        bullets = []
        for line in response.content.strip().split('\n'):
            line = line.strip()
            if line.startswith('- '):
                bullet = line[2:].strip()
                if len(bullet) > 120:
                    bullet = bullet[:117] + "..."
                bullets.append(bullet)

        while len(bullets) < 3:
            bullets.append("Additional research pending")

        return bullets[:3]

    def extract_bullets(self, content: str, section_title: str) -> List[str]:
        """Extract bullets from markdown or condense if needed."""
        bullets = []
        for line in content.split('\n'):
            line = line.strip()
            if line.startswith('- '):
                bullet = line[2:].strip()
                if bullet:
                    bullets.append(bullet)

        if len(bullets) == 3:
            return bullets

        return self.condense_to_three_bullets(content, section_title)

    def create_presentation(self, strategy_data: Dict[str, Any]) -> Presentation:
        """Create a PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        sector = strategy_data.get('sector', 'Networking Strategy')
        region = strategy_data.get('region', 'US')

        print("  Creating slides...")

        # Title Slide
        print("    -> Title slide")
        self._add_title_slide(prs, sector, region)

        # Content Slides
        for section in strategy_data.get('sections', []):
            section_num = section.get('number', '?')
            section_title = section.get('title', 'Section')
            full_title = f"{section_num}. {section_title}"

            print(f"    -> Section {section_num}: {section_title}")
            bullets = self.extract_bullets(section['content'], section_title)
            self._add_content_slide(prs, full_title, bullets)

        # Closing Slide
        print("    -> Closing slide")
        self._add_closing_slide(prs, sector)

        return prs

    def _add_title_slide(self, prs: Presentation, sector: str, region: str):
        """Create title slide using the same format as VC thesis deck."""
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.style.PRIMARY
        bar.line.fill.background()

        title_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(2.5), Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Networking Strategy"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.style.PRIMARY

        sub_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(3.7), Inches(12), Inches(0.8)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sector
        p.font.size = Pt(32)
        p.font.color.rgb = self.style.SECONDARY

        meta_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(5.2), Inches(12), Inches(0.5)
        )
        tf = meta_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{region}  |  {datetime.now().strftime('%B %Y')}  |  GPTOSS VC Thesis System"
        p.font.size = Pt(14)
        p.font.color.rgb = self.style.TEXT_LIGHT

        bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = self.style.ACCENT
        bar2.line.fill.background()

    def _add_content_slide(self, prs: Presentation, title: str, bullets: List[str]):
        """Create a content slide with title and exactly 3 bullets."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.style.PRIMARY

        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()

        for i, bullet_text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet_text
            p.font.size = Pt(20)
            p.font.color.rgb = self.style.TEXT_DARK
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            p.level = 0

    def _add_closing_slide(self, prs: Presentation, sector: str):
        """Create professional closing slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.style.PRIMARY
        bg.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.style.WHITE
        p.alignment = PP_ALIGN.CENTER

        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{sector} Networking Strategy"
        p.font.size = Pt(20)
        p.font.color.rgb = self.style.ACCENT
        p.alignment = PP_ALIGN.CENTER

        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.3), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated by GPTOSS VC Thesis System  |  {datetime.now().strftime('%B %Y')}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER

    def save_presentation(self, prs: Presentation, sector: str) -> Path:
        """Save the PowerPoint presentation."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        sector_slug = sector.lower().replace(" ", "_")
        filename = f"networking_strategy_{sector_slug}_{timestamp}.pptx"

        output_path = self.outputs_path / filename
        prs.save(str(output_path))

        print(f"\n  [OK] PowerPoint saved to: {output_path}")
        return output_path

    def run(self, sector: str, strategy_file: Path = None) -> Path:
        """Run the complete PowerPoint generation workflow."""
        print(f"\n{'='*60}")
        print("POWERPOINT_NETWORKING_AGENT: Creating Presentation")
        print(f"Sector: {sector}")
        print(f"{'='*60}\n")

        print("  Loading strategy...")
        strategy_content = self.load_strategy(strategy_file=strategy_file, sector=sector)

        print("  Parsing strategy structure...")
        strategy_data = self.parse_strategy(strategy_content)
        strategy_data['sector'] = sector

        prs = self.create_presentation(strategy_data)
        output_path = self.save_presentation(prs, sector)

        slide_count = len(prs.slides)
        print(f"\n{'='*60}")
        print("[OK] PowerPoint generation complete!")
        print(f"     {slide_count} slides | 3 bullets each")
        print(f"{'='*60}\n")

        return output_path


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        print("Usage: python agent.py <sector> [strategy_file]")
        print("Example: python agent.py 'B2B Fintech'")
        sys.exit(1)

    sector = sys.argv[1]
    strategy_file = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    try:
        agent = PowerPointNetworkingAgent()
        agent.run(sector, strategy_file=strategy_file)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
