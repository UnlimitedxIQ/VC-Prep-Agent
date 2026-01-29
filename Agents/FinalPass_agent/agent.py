"""
FinalPass_agent: PowerPoint Quality Assurance

Validates and fixes PowerPoint presentations:
- Ensures all fields are filled in
- Each content slide has exactly 3 key bullets
- Grammar and clarity checks via LLM
- Text overflow and spacing validation
- Generates corrected presentation
"""

import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Tuple
import re
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

env_path = Path(__file__).parent.parent.parent / "Secrets" / ".env"
load_dotenv(env_path)


class QAIssue:
    """Represents a quality issue found in the presentation."""
    def __init__(self, slide_num: int, issue_type: str, description: str, severity: str = "warning"):
        self.slide_num = slide_num
        self.issue_type = issue_type
        self.description = description
        self.severity = severity  # "error", "warning", "info"
        self.fixed = False


class FinalPassAgent:
    """Agent for validating and fixing PowerPoint presentations."""

    def __init__(self, api_key: str = None):
        """Initialize the FinalPass Agent."""
        self.openai_api_key = api_key or os.getenv("OPENAI_API_KEY")

        if not self.openai_api_key:
            raise ValueError("OpenAI API key not provided")

        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.2,
            openai_api_key=self.openai_api_key
        )

        self.base_path = Path(__file__).parent.parent.parent
        self.outputs_path = self.base_path / "Outputs"
        self.outputs_path.mkdir(exist_ok=True)

        self.issues: List[QAIssue] = []

        # Layout constants
        self.MAX_CHARS_PER_LINE = 85
        self.MAX_LINES_PER_TEXTBOX = 8
        self.TARGET_BULLETS = 3

    def load_presentation(self, pptx_file: Path = None, sector: str = None) -> Tuple[Path, Presentation]:
        """Load a PowerPoint presentation."""
        if pptx_file and pptx_file.exists():
            prs = Presentation(str(pptx_file))
            return pptx_file, prs

        if sector:
            sector_slug = sector.lower().replace(" ", "_")
            pattern = f"vc_thesis_{sector_slug}_*.pptx"
            matching_files = [f for f in self.outputs_path.glob(pattern) if "_reviewed" not in f.name]

            if matching_files:
                latest_file = max(matching_files, key=lambda p: p.stat().st_mtime)
                print(f"  Loading: {latest_file.name}")
                prs = Presentation(str(latest_file))
                return latest_file, prs

        raise ValueError("No presentation file found.")

    def extract_slide_content(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Extract content from all slides for analysis."""
        slides_data = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                'number': i + 1,
                'title': '',
                'text_boxes': [],
                'total_text': '',
                'bullet_count': 0,
                'shapes': []
            }

            for shape in slide.shapes:
                shape_info = {
                    'type': type(shape).__name__,
                    'has_text': hasattr(shape, "text_frame"),
                    'text': '',
                    'width': shape.width,
                    'height': shape.height,
                    'left': shape.left,
                    'top': shape.top
                }

                if hasattr(shape, "text_frame"):
                    text = ""
                    bullet_count = 0
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            text += para_text + "\n"
                            # Count as bullet if it's a content paragraph
                            if len(para_text) > 10:
                                bullet_count += 1

                    shape_info['text'] = text.strip()
                    shape_info['bullet_count'] = bullet_count

                    if text.strip():
                        slide_data['text_boxes'].append(shape_info)
                        slide_data['total_text'] += text
                        slide_data['bullet_count'] += bullet_count

                        # Check if this is the title
                        if shape == slide.shapes.title:
                            slide_data['title'] = text.strip().split('\n')[0]

                slide_data['shapes'].append(shape_info)

            slides_data.append(slide_data)

        return slides_data

    def validate_content_completeness(self, slides_data: List[Dict[str, Any]]):
        """Check that all slides have required content."""
        print("  Checking content completeness...")

        for slide in slides_data:
            num = slide['number']

            # Skip title (1) and closing (last) slides
            if num == 1 or num == len(slides_data):
                continue

            # Check for empty slides
            if not slide['total_text'].strip():
                self.issues.append(QAIssue(
                    num, "empty_slide",
                    "Slide has no text content",
                    "error"
                ))

            # Check for missing title
            if not slide['title']:
                self.issues.append(QAIssue(
                    num, "missing_title",
                    "Slide is missing a title",
                    "error"
                ))

    def validate_bullet_count(self, slides_data: List[Dict[str, Any]]):
        """Ensure content slides have appropriate bullet counts."""
        print("  Checking bullet counts...")

        # Content slides are typically 3-7 (excluding title and closing)
        for slide in slides_data:
            num = slide['number']

            # Skip title, exec summary, and closing slides
            if num in [1, 2, len(slides_data)]:
                continue

            bullet_count = slide['bullet_count']

            if bullet_count < 2:
                self.issues.append(QAIssue(
                    num, "too_few_bullets",
                    f"Slide has only {bullet_count} bullet(s), recommend at least 3",
                    "warning"
                ))
            elif bullet_count > 6:
                self.issues.append(QAIssue(
                    num, "too_many_bullets",
                    f"Slide has {bullet_count} bullets, may be overcrowded",
                    "warning"
                ))

    def validate_text_overflow(self, slides_data: List[Dict[str, Any]]):
        """Check for potential text overflow issues."""
        print("  Checking text overflow...")

        for slide in slides_data:
            num = slide['number']

            for tb in slide['text_boxes']:
                text = tb['text']
                lines = text.split('\n')

                # Check for very long lines
                for line in lines:
                    if len(line) > self.MAX_CHARS_PER_LINE * 1.5:
                        self.issues.append(QAIssue(
                            num, "long_line",
                            f"Line exceeds recommended length ({len(line)} chars): '{line[:50]}...'",
                            "warning"
                        ))

                # Check for too many lines in a text box
                if len(lines) > self.MAX_LINES_PER_TEXTBOX:
                    self.issues.append(QAIssue(
                        num, "too_many_lines",
                        f"Text box has {len(lines)} lines, may overflow",
                        "warning"
                    ))

    def check_grammar_and_clarity(self, slides_data: List[Dict[str, Any]]) -> Dict[int, str]:
        """Use LLM to check grammar and suggest improvements."""
        print("  Checking grammar and clarity...")

        # Combine all slide text for efficient LLM call
        all_text = ""
        for slide in slides_data:
            if slide['total_text'].strip():
                all_text += f"\n\n--- SLIDE {slide['number']} ---\n{slide['total_text']}"

        system_prompt = """You are a professional editor reviewing presentation text.
For each slide, identify:
1. Grammar errors (spelling, punctuation, sentence structure)
2. Unclear or vague statements
3. Inconsistent formatting or style

Return ONLY a JSON-like format with fixes needed:
SLIDE X: [issue] -> [suggested fix]

If a slide is fine, write: SLIDE X: OK

Be concise and specific."""

        user_prompt = f"""Review this presentation text for grammar and clarity issues:
{all_text}

List specific issues and fixes for each slide."""

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]

        response = self.llm.invoke(messages)
        grammar_feedback = response.content

        # Parse feedback and add issues
        for line in grammar_feedback.split('\n'):
            if 'SLIDE' in line and 'OK' not in line and '->' in line:
                try:
                    slide_match = re.search(r'SLIDE\s*(\d+)', line)
                    if slide_match:
                        slide_num = int(slide_match.group(1))
                        issue_text = line.split(':', 1)[1].strip() if ':' in line else line
                        self.issues.append(QAIssue(
                            slide_num, "grammar",
                            issue_text,
                            "info"
                        ))
                except:
                    pass

        return grammar_feedback

    def generate_fixes(self, slides_data: List[Dict[str, Any]]) -> Dict[int, Dict[str, str]]:
        """Generate LLM-based content fixes for problem slides."""
        print("  Generating content fixes...")

        fixes = {}

        # Find slides that need bullet consolidation
        problem_slides = []
        for slide in slides_data:
            num = slide['number']
            if num in [1, 2, len(slides_data)]:
                continue

            if slide['bullet_count'] < 2 or slide['bullet_count'] > 6:
                problem_slides.append(slide)

        if not problem_slides:
            print("    No major content fixes needed")
            return fixes

        for slide in problem_slides:
            num = slide['number']
            text = slide['total_text']

            system_prompt = """You are a presentation editor. Rewrite the given slide content to have EXACTLY 3 clear, concise bullet points.

Rules:
- Each bullet should be 1-2 sentences max
- Keep the key information and data points
- Maintain any citations/sources
- Use parallel structure
- Be specific and actionable

Return ONLY the 3 bullets, one per line, starting with "- " """

            user_prompt = f"""Rewrite this slide content into exactly 3 clear bullets:

{text}

Return only the 3 bullets."""

            messages = [
                SystemMessage(content=system_prompt),
                HumanMessage(content=user_prompt)
            ]

            response = self.llm.invoke(messages)
            fixes[num] = {
                'original': text,
                'fixed': response.content.strip()
            }

        return fixes

    def apply_fixes_to_presentation(self, prs: Presentation, fixes: Dict[int, Dict[str, str]]) -> Presentation:
        """Apply generated fixes to the presentation."""
        print("  Applying fixes to presentation...")

        # Note: Full programmatic editing of existing slides is complex
        # For now, we'll add the fixes as speaker notes for manual review
        # A more advanced version could rebuild slides

        for slide_num, fix_data in fixes.items():
            if slide_num <= len(prs.slides):
                slide = prs.slides[slide_num - 1]

                # Add fix suggestion to notes
                notes_text = f"""SUGGESTED FIX (3 bullets):

{fix_data['fixed']}

---
Original had issues with bullet count.
Review and apply manually if needed."""

                if not slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                notes_frame = slide.notes_slide.notes_text_frame
                notes_frame.text = notes_text

        return prs

    def generate_qa_report(self, slides_data: List[Dict[str, Any]], grammar_feedback: str) -> str:
        """Generate a comprehensive QA report."""
        report = f"""# PowerPoint Quality Assurance Report

**Generated**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
**Slides Analyzed**: {len(slides_data)}
**Issues Found**: {len(self.issues)}

---

## Summary

| Severity | Count |
|----------|-------|
| Errors   | {len([i for i in self.issues if i.severity == 'error'])} |
| Warnings | {len([i for i in self.issues if i.severity == 'warning'])} |
| Info     | {len([i for i in self.issues if i.severity == 'info'])} |

---

## Issues by Slide

"""
        # Group issues by slide
        slides_with_issues = {}
        for issue in self.issues:
            if issue.slide_num not in slides_with_issues:
                slides_with_issues[issue.slide_num] = []
            slides_with_issues[issue.slide_num].append(issue)

        if not slides_with_issues:
            report += "**No issues found! Presentation passes QA.**\n\n"
        else:
            for slide_num in sorted(slides_with_issues.keys()):
                issues = slides_with_issues[slide_num]
                report += f"### Slide {slide_num}\n\n"
                for issue in issues:
                    icon = "X" if issue.severity == "error" else ("!" if issue.severity == "warning" else "i")
                    report += f"- [{icon}] **{issue.issue_type}**: {issue.description}\n"
                report += "\n"

        report += """---

## Grammar & Clarity Analysis

"""
        report += grammar_feedback

        report += """

---

## Checklist

- [ ] All slides have titles
- [ ] Content slides have 3 key bullets each
- [ ] No text overflow or wrapping issues
- [ ] Grammar and spelling verified
- [ ] Citations included where needed
- [ ] Consistent formatting throughout

---

*Report generated by FinalPass_agent*
"""
        return report

    def save_presentation(self, prs: Presentation, sector: str) -> Path:
        """Save the reviewed presentation."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        sector_slug = sector.lower().replace(" ", "_")
        filename = f"vc_thesis_{sector_slug}_{timestamp}_reviewed.pptx"

        output_path = self.outputs_path / filename
        prs.save(str(output_path))

        print(f"\n  [OK] Reviewed presentation saved to: {output_path}")
        return output_path

    def save_report(self, report: str, sector: str) -> Path:
        """Save the QA report."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        sector_slug = sector.lower().replace(" ", "_")
        filename = f"qa_report_{sector_slug}_{timestamp}.md"

        output_path = self.outputs_path / filename
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(report)

        print(f"  [OK] QA report saved to: {output_path}")
        return output_path

    def run(self, sector: str, pptx_file: Path = None) -> Tuple[Path, Path]:
        """Run the complete QA workflow."""
        print(f"\n{'='*60}")
        print(f"FinalPass_agent: Quality Assurance Review")
        print(f"Sector: {sector}")
        print(f"{'='*60}\n")

        # Reset issues
        self.issues = []

        # Load presentation
        print("Loading presentation...")
        original_path, prs = self.load_presentation(pptx_file=pptx_file, sector=sector)
        print(f"  Found {len(prs.slides)} slides\n")

        # Extract content
        print("Analyzing slides...")
        slides_data = self.extract_slide_content(prs)

        # Run validations
        print("\nRunning validations...")
        self.validate_content_completeness(slides_data)
        self.validate_bullet_count(slides_data)
        self.validate_text_overflow(slides_data)

        # Grammar check
        print("\nRunning grammar check...")
        grammar_feedback = self.check_grammar_and_clarity(slides_data)

        # Generate and apply fixes
        print("\nGenerating fixes...")
        fixes = self.generate_fixes(slides_data)

        if fixes:
            prs = self.apply_fixes_to_presentation(prs, fixes)

        # Generate report
        print("\nGenerating QA report...")
        report = self.generate_qa_report(slides_data, grammar_feedback)

        # Save outputs
        print("\nSaving outputs...")
        reviewed_path = self.save_presentation(prs, sector)
        report_path = self.save_report(report, sector)

        # Print summary
        error_count = len([i for i in self.issues if i.severity == 'error'])
        warning_count = len([i for i in self.issues if i.severity == 'warning'])

        print(f"\n{'='*60}")
        print(f"[OK] QA Review Complete!")
        print(f"     Errors: {error_count} | Warnings: {warning_count}")
        print(f"{'='*60}\n")

        return reviewed_path, report_path


def main():
    """CLI entry point for FinalPass_agent."""
    if len(sys.argv) < 2:
        print("Usage: python agent.py <sector> [pptx_file]")
        print("Example: python agent.py 'B2B Fintech'")
        sys.exit(1)

    sector = sys.argv[1]
    pptx_file = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    try:
        agent = FinalPassAgent()
        agent.run(sector, pptx_file=pptx_file)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
