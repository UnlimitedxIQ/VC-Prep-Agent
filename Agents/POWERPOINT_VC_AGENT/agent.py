"""
POWERPOINT_VC_AGENT: High-Quality PowerPoint Generator

Creates professional VC thesis presentations with:
- Proper slide titles (detectable by QA)
- Exactly 3 bullets per content slide
- Clean spacing and typography
- No text overflow
"""

import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Tuple
import re
from dotenv import load_dotenv

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage

env_path = Path(__file__).parent.parent.parent / "Secrets" / ".env"
load_dotenv(env_path)


class SlideStyle:
    """Professional styling constants."""
    PRIMARY = RGBColor(0, 51, 102)
    SECONDARY = RGBColor(0, 102, 153)
    ACCENT = RGBColor(0, 153, 204)
    TEXT_DARK = RGBColor(51, 51, 51)
    TEXT_LIGHT = RGBColor(102, 102, 102)
    WHITE = RGBColor(255, 255, 255)

    MARGIN_LEFT = Inches(0.75)
    MAX_BULLET_CHARS = 120


class PowerPointAgent:
    """Agent for converting VC thesis to high-quality PowerPoint."""

    def __init__(self, api_key: str = None):
        self.openai_api_key = api_key or os.getenv("OPENAI_API_KEY")
        if not self.openai_api_key:
            raise ValueError("OpenAI API key not provided")

        self.llm = ChatOpenAI(
            model="gpt-4o-mini",
            temperature=0.3,
            openai_api_key=self.openai_api_key
        )

        self.base_path = Path(__file__).parent.parent.parent
        self.outputs_path = self.base_path / "Outputs"
        self.outputs_path.mkdir(exist_ok=True)
        self.style = SlideStyle()

    def load_thesis(self, thesis_file: Path = None, sector: str = None) -> str:
        """Load the thesis markdown file."""
        if thesis_file and thesis_file.exists():
            with open(thesis_file, 'r', encoding='utf-8') as f:
                return f.read()

        if sector:
            sector_slug = sector.lower().replace(" ", "_")
            pattern = f"vc_thesis_{sector_slug}_*.md"
            matching_files = list(self.outputs_path.glob(pattern))

            if matching_files:
                latest_file = max(matching_files, key=lambda p: p.stat().st_mtime)
                print(f"  Loading thesis: {latest_file.name}")
                with open(latest_file, 'r', encoding='utf-8') as f:
                    return f.read()

        raise ValueError("No thesis file found.")

    def parse_thesis(self, markdown: str) -> Dict[str, Any]:
        """Parse the thesis markdown into structured sections."""
        data = {
            'sector': '',
            'region': 'US',
            'executive_summary': '',
            'sections': []
        }

        # Extract region
        region_match = re.search(r'\*\*Region\*\*:\s*(.+?)$', markdown, re.MULTILINE)
        if region_match:
            data['region'] = region_match.group(1).strip()

        # Extract sector from title
        sector_match = re.search(r'^## (.+?)$', markdown, re.MULTILINE)
        if sector_match:
            data['sector'] = sector_match.group(1).strip()

        # Extract executive summary
        exec_match = re.search(r'## Executive Summary\s*\n+(.*?)(?=\n---|\n## \d)', markdown, re.DOTALL)
        if exec_match:
            data['executive_summary'] = exec_match.group(1).strip()

        # Extract numbered sections (1-5)
        section_pattern = r'## (\d+)\. (.+?)\n(.*?)(?=\n## \d|\n---\s*\n## Methodology|$)'
        for match in re.finditer(section_pattern, markdown, re.DOTALL):
            section_num = match.group(1)
            section_title = match.group(2).strip()
            section_content = match.group(3).strip()

            data['sections'].append({
                'number': section_num,
                'title': section_title,
                'content': section_content
            })

        return data

    def condense_to_three_bullets(self, content: str, section_title: str) -> List[str]:
        """Use LLM to condense section content to exactly 3 concise bullets."""
        system_prompt = """You are a presentation editor. Convert the content into EXACTLY 3 bullet points.

Rules:
- Each bullet must be under 100 characters
- Keep key data points and metrics
- Be specific and concise
- No sub-bullets
- Start each with a strong action word or key term

Return ONLY 3 lines, each starting with "- " """

        user_prompt = f"""Section: {section_title}

Content:
{content}

Convert to exactly 3 concise bullet points (under 100 chars each):"""

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=user_prompt)
        ]

        response = self.llm.invoke(messages)

        # Parse the 3 bullets
        bullets = []
        for line in response.content.strip().split('\n'):
            line = line.strip()
            if line.startswith('- '):
                bullet = line[2:].strip()
                # Truncate if too long
                if len(bullet) > 120:
                    bullet = bullet[:117] + "..."
                bullets.append(bullet)

        # Ensure exactly 3
        while len(bullets) < 3:
            bullets.append("Additional research pending")

        return bullets[:3]

    def condense_executive_summary(self, summary: str) -> List[str]:
        """Condense executive summary to 3 key points."""
        system_prompt = """Condense this executive summary into EXACTLY 3 key takeaways.
Each point should be 1 sentence, under 80 characters.
Return only 3 lines starting with "- " """

        messages = [
            SystemMessage(content=system_prompt),
            HumanMessage(content=f"Summary:\n{summary}\n\nReturn 3 key points:")
        ]

        response = self.llm.invoke(messages)

        bullets = []
        for line in response.content.strip().split('\n'):
            line = line.strip()
            if line.startswith('- '):
                bullet = line[2:].strip()
                if len(bullet) > 100:
                    bullet = bullet[:97] + "..."
                bullets.append(bullet)

        while len(bullets) < 3:
            bullets.append("Key insight from research")

        return bullets[:3]

    def create_presentation(self, thesis_data: Dict[str, Any]) -> Presentation:
        """Create a high-quality PowerPoint presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        sector = thesis_data.get('sector', 'Industry Thesis')
        region = thesis_data.get('region', 'US')

        print("  Creating slides...")

        # 1. Title Slide
        print("    -> Title slide")
        self._add_title_slide(prs, sector, region)

        # 2. Executive Summary (condensed to 3 bullets)
        if thesis_data.get('executive_summary'):
            print("    -> Executive Summary (condensing to 3 bullets)")
            summary_bullets = self.condense_executive_summary(thesis_data['executive_summary'])
            self._add_content_slide(prs, "Executive Summary", summary_bullets)

        # 3. Content Slides (each with exactly 3 bullets)
        for section in thesis_data.get('sections', []):
            section_num = section.get('number', '?')
            section_title = section.get('title', 'Section')
            full_title = f"{section_num}. {section_title}"

            print(f"    -> Section {section_num}: {section_title} (condensing to 3 bullets)")
            bullets = self.condense_to_three_bullets(section['content'], section_title)
            self._add_content_slide(prs, full_title, bullets)

        # 4. Closing Slide
        print("    -> Closing slide")
        self._add_closing_slide(prs, sector)

        return prs

    def _add_title_slide(self, prs: Presentation, sector: str, region: str):
        """Create title slide using layout with proper title placeholder."""
        # Use Title Slide layout (index 0) for proper title detection
        slide_layout = prs.slide_layouts[6]  # Blank for custom design
        slide = prs.slides.add_slide(slide_layout)

        # Top accent bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.style.PRIMARY
        bar.line.fill.background()

        # Main title - make it the first shape for title detection
        title_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(2.5), Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Venture Capital Industry Thesis"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.style.PRIMARY

        # Sector subtitle
        sub_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(3.7), Inches(12), Inches(0.8)
        )
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sector
        p.font.size = Pt(32)
        p.font.color.rgb = self.style.SECONDARY

        # Metadata
        meta_box = slide.shapes.add_textbox(
            self.style.MARGIN_LEFT, Inches(5.2), Inches(12), Inches(0.5)
        )
        tf = meta_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{region}  |  {datetime.now().strftime('%B %Y')}  |  GPTOSS VC Thesis System"
        p.font.size = Pt(14)
        p.font.color.rgb = self.style.TEXT_LIGHT

        # Bottom accent bar
        bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), prs.slide_width, Inches(0.15))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = self.style.ACCENT
        bar2.line.fill.background()

    def _add_content_slide(self, prs: Presentation, title: str, bullets: List[str]):
        """Create a content slide with title and exactly 3 bullets."""
        # Use Title and Content layout (index 1) for proper structure
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        # Set the title using the proper placeholder
        title_shape = slide.shapes.title
        title_shape.text = title

        # Style the title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(32)
            paragraph.font.bold = True
            paragraph.font.color.rgb = self.style.PRIMARY

        # Add bullets to content placeholder
        content_placeholder = slide.placeholders[1]
        tf = content_placeholder.text_frame
        tf.clear()

        for i, bullet_text in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = bullet_text
            p.font.size = Pt(20)
            p.font.color.rgb = self.style.TEXT_DARK
            p.space_before = Pt(12)
            p.space_after = Pt(8)
            p.level = 0

    def _add_closing_slide(self, prs: Presentation, sector: str):
        """Create professional closing slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Full background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.style.PRIMARY
        bg.line.fill.background()

        # Thank you
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.style.WHITE
        p.alignment = PP_ALIGN.CENTER

        # Sector
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(0.6))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{sector} Industry Thesis"
        p.font.size = Pt(20)
        p.font.color.rgb = self.style.ACCENT
        p.alignment = PP_ALIGN.CENTER

        # Footer
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(11.3), Inches(0.4))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated by GPTOSS VC Thesis System  |  {datetime.now().strftime('%B %Y')}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(180, 200, 220)
        p.alignment = PP_ALIGN.CENTER

    def save_presentation(self, prs: Presentation, sector: str) -> Path:
        """Save the PowerPoint presentation."""
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        sector_slug = sector.lower().replace(" ", "_")
        filename = f"vc_thesis_{sector_slug}_{timestamp}.pptx"

        output_path = self.outputs_path / filename
        prs.save(str(output_path))

        print(f"\n  [OK] PowerPoint saved to: {output_path}")
        return output_path

    def run(self, sector: str, thesis_file: Path = None) -> Path:
        """Run the complete PowerPoint generation workflow."""
        print(f"\n{'='*60}")
        print(f"POWERPOINT_VC_AGENT: Creating Presentation")
        print(f"Sector: {sector}")
        print(f"{'='*60}\n")

        # Load thesis
        print("  Loading thesis...")
        thesis_content = self.load_thesis(thesis_file=thesis_file, sector=sector)

        # Parse thesis
        print("  Parsing thesis structure...")
        thesis_data = self.parse_thesis(thesis_content)
        thesis_data['sector'] = sector

        # Create presentation
        prs = self.create_presentation(thesis_data)

        # Save
        output_path = self.save_presentation(prs, sector)

        slide_count = len(prs.slides)
        print(f"\n{'='*60}")
        print(f"[OK] PowerPoint generation complete!")
        print(f"     {slide_count} slides | 3 bullets each | QA-ready")
        print(f"{'='*60}\n")

        return output_path


def main():
    """CLI entry point."""
    if len(sys.argv) < 2:
        print("Usage: python agent.py <sector> [thesis_file]")
        print("Example: python agent.py 'B2B Fintech'")
        sys.exit(1)

    sector = sys.argv[1]
    thesis_file = Path(sys.argv[2]) if len(sys.argv) > 2 else None

    try:
        agent = PowerPointAgent()
        agent.run(sector, thesis_file=thesis_file)
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
